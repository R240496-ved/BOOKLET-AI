"""
Input Extraction Module (Enhanced)
Extracts text AND images from PDF, PPTX, and TXT files.
Text extraction is optimised for clean RAG-quality output.
Images are INTERLEAVED with text using coordinate sorting.

Image metadata format:
  {
    "page": int,        # page number (PDF) or slide index (PPTX)
    "source": str,      # "pdf" or "pptx"
    "path": str,        # absolute path to saved image file
    "name": str,        # filename
  }
"""

import os
import re
import unicodedata
from collections import Counter


# ─────────────────────────────────────────────
# MAIN ENTRY POINT
# ─────────────────────────────────────────────

def extract(file_path: str, image_output_dir: str = "data/extracted_images") -> dict:
    """
    Extract clean text with inline image markers.
    Returns: {"text": str, "images": list[dict]}
    """
    os.makedirs(image_output_dir, exist_ok=True)
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path, image_output_dir)
    elif ext == ".pptx":
        return _extract_pptx(file_path, image_output_dir)
    elif ext == ".txt":
        return _extract_txt(file_path)
    elif ext in [".png", ".jpg", ".jpeg"]:
        return _extract_img(file_path, image_output_dir)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


# ─────────────────────────────────────────────
# TEXT POST-PROCESSING UTILITIES
# ─────────────────────────────────────────────

_LIGATURE_MAP = {
    "\ufb00": "ff", "\ufb01": "fi", "\ufb02": "fl",
    "\ufb03": "ffi", "\ufb04": "ffl", "\ufb05": "st",
    "\u2019": "'", "\u2018": "'", "\u201c": '"', "\u201d": '"',
    "\u2013": "-", "\u2014": "--", "\u00a0": " ",
}

def _fix_ligatures(text: str) -> str:
    for bad, good in _LIGATURE_MAP.items():
        text = text.replace(bad, good)
    return text


def _is_page_number(line: str) -> bool:
    """True if the line is just a page number."""
    stripped = line.strip()
    return bool(re.fullmatch(r"(page\s*)?\d+(\s*of\s*\d+)?", stripped, re.IGNORECASE)
                or re.fullmatch(r"[-–]\s*\d+\s*[-–]", stripped))


def _find_repeating_lines(pages_text: list[str], min_freq_ratio: float = 0.4) -> set[str]:
    """Detect header/footer lines that repeat across many pages."""
    total_pages = len(pages_text)
    if total_pages < 3:
        return set()

    line_counts: Counter = Counter()

    for page_text in pages_text:
        lines = page_text.split("\n")
        candidates = set(lines[:2] + lines[-2:])
        for c in candidates:
            c = c.strip()
            if c and len(c) > 3:
                line_counts[c] += 1

    threshold = total_pages * min_freq_ratio
    return {line for line, count in line_counts.items() if count >= threshold}


def _clean_page_content(items: list[dict], boilerplate: set[str]) -> str:
    """
    Cleans items and joins them into a single text block.
    Items are either {'type': 'text', 'content': str} or {'type': 'image', 'path': str}.
    """
    kept_lines = []

    for item in items:
        if item['type'] == 'image':
            kept_lines.append(f"\n![IMAGE:{item['path']}]\n")
            continue

        raw = _fix_ligatures(item['content'])
        raw = unicodedata.normalize("NFKC", raw)

        lines = raw.split("\n")
        for line in lines:
            stripped = line.strip()
            if not stripped or stripped in boilerplate or _is_page_number(stripped):
                continue
            kept_lines.append(stripped)

    # Merge hyphenated line-breaks
    merged = []
    i = 0
    while i < len(kept_lines):
        line = kept_lines[i]
        if line.endswith("-") and i + 1 < len(kept_lines) and not kept_lines[i+1].startswith("![IMAGE:"):
            merged.append(line[:-1] + kept_lines[i + 1])
            i += 2
        else:
            merged.append(line)
            i += 1

    return "\n".join(merged).strip()


# ─────────────────────────────────────────────
# PDF EXTRACTION
# ─────────────────────────────────────────────

def _extract_pdf(file_path: str, image_output_dir: str) -> dict:
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    page_data_list = []  # List of list of items (text or image) per page
    images_metadata = []

    for page_num, page in enumerate(doc):
        page_items = []

        # 1. Get Text Blocks
        text_blocks = page.get_text("blocks")
        for b in text_blocks:
            # (x0, y0, x1, y1, "text", block_no, block_type)
            if b[6] == 0:  # Text
                page_items.append({
                    "type": "text",
                    "content": b[4],
                    "y": b[1],
                    "x": b[0]
                })

        # 2. Get Images
        image_list = page.get_images(full=True)
        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            # Get image position
            rects = page.get_image_rects(xref)
            if not rects:
                continue
            rect = rects[0] # Take first occurrence

            try:
                base_image = doc.extract_image(xref)
                if base_image["width"] < 60 or base_image["height"] < 60:
                    continue

                img_ext = base_image["ext"]
                img_name = f"pdf_p{page_num + 1}_img{img_idx + 1}.{img_ext}"
                img_path = os.path.join(image_output_dir, img_name)
                
                with open(img_path, "wb") as f:
                    f.write(base_image["image"])

                img_meta = {
                    "page": page_num + 1,
                    "source": "pdf",
                    "path": img_path,
                    "name": img_name,
                }
                images_metadata.append(img_meta)
                
                page_items.append({
                    "type": "image",
                    "path": img_path,
                    "y": rect.y0,
                    "x": rect.x0
                })
            except Exception as e:
                print(f"[Extraction] Skipped image on page {page_num + 1}: {e}")

        # Sort items by Y, then X (reading order)
        page_items.sort(key=lambda item: (item['y'], item['x']))
        page_data_list.append(page_items)

    doc.close()

    # Detect repeating lines across simple text extraction to find boilerplate
    flat_text_per_page = []
    for items in page_data_list:
        flat_text_per_page.append("\n".join([it['content'] for it in items if it['type'] == 'text']))
    
    boilerplate = _find_repeating_lines(flat_text_per_page)

    # Clean and combine
    final_text_blocks = []
    for page_items in page_data_list:
        cleaned = _clean_page_content(page_items, boilerplate)
        if cleaned:
            final_text_blocks.append(cleaned)

    return {
        "text": "\n\n".join(final_text_blocks),
        "images": images_metadata,
    }


# ─────────────────────────────────────────────
# PPTX EXTRACTION
# ─────────────────────────────────────────────

def _extract_pptx(file_path: str, image_output_dir: str) -> dict:
    from pptx import Presentation

    prs = Presentation(file_path)
    slide_data_list = []
    images_metadata = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_items = []

        for shape in slide.shapes:
            # Get position
            y = shape.top if hasattr(shape, "top") else 0
            x = shape.left if hasattr(shape, "left") else 0

            # Text
            if hasattr(shape, "text_frame"):
                full_text = []
                for para in shape.text_frame.paragraphs:
                    para_text = " ".join(run.text for run in para.runs if run.text.strip()).strip()
                    if para_text:
                        full_text.append(para_text)
                if full_text:
                    slide_items.append({
                        "type": "text",
                        "content": "\n".join(full_text),
                        "y": y,
                        "x": x
                    })
            elif hasattr(shape, "text") and shape.text.strip():
                slide_items.append({
                    "type": "text",
                    "content": shape.text.strip(),
                    "y": y,
                    "x": x
                })

            # Get Images (Pictures and Grouped Pictures)
            def _extract_shape_images(s, s_idx, depth=0):
                if s.shape_type == 13: # PICTURE
                    try:
                        image = s.image
                        img_ext = image.ext
                        img_name = f"pptx_s{s_idx + 1}_d{depth}_img{len(images_metadata) + 1}.{img_ext}"
                        img_path = os.path.join(image_output_dir, img_name)
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        
                        meta = {"page": s_idx + 1, "source": "pptx", "path": img_path, "name": img_name}
                        images_metadata.append(meta)
                        slide_items.append({"type": "image", "path": img_path, "y": y, "x": x})
                    except: pass
                elif s.shape_type == 6: # GROUP
                    for subshape in s.shapes:
                        _extract_shape_images(subshape, s_idx, depth + 1)

            _extract_shape_images(shape, slide_idx)

        # Sort by Y then X
        slide_items.sort(key=lambda item: (item['y'], item['x']))
        slide_data_list.append(slide_items)

    # Clean and combine
    final_text_blocks = []
    for slide_items in slide_data_list:
        cleaned = _clean_page_content(slide_items, set()) # No boilerplate check for PPT
        if cleaned:
            final_text_blocks.append(cleaned)

    return {
        "text": "\n\n".join(final_text_blocks),
        "images": images_metadata,
    }


# ─────────────────────────────────────────────
# TXT EXTRACTION
# ─────────────────────────────────────────────

def _extract_txt(file_path: str) -> dict:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return {
            "text": _fix_ligatures(f.read().strip()),
            "images": [],
        }


# ─────────────────────────────────────────────
# LEGACY COMPAT
# ─────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    return extract(file_path)["text"]

def _extract_img(file_path: str, image_output_dir: str) -> dict:
    import shutil
    import os
    from modules.ocr_module import extract_text_from_image

    filename = os.path.basename(file_path)
    img_dest_path = os.path.join(image_output_dir, filename)
    if os.path.abspath(file_path) != os.path.abspath(img_dest_path):
        shutil.copy2(file_path, img_dest_path)

    image_metadata = {
        "page": 1,
        "source": "image",
        "path": img_dest_path,
        "name": filename
    }

    # Use the optimized OCR module extraction
    cleaned_ocr = extract_text_from_image(img_dest_path)

    final_text = f"![IMAGE:{img_dest_path}]\n\n{cleaned_ocr}"

    return {
        "text": final_text,
        "images": [image_metadata],
    }