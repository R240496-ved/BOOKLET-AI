"""
PPT Generator Module (Specialized Visual Edition)
Generates a structured .pptx presentation from Markdown-structured text.
Strategy: Dedicated Visual Slides for images, Dedicated Content Slides for text.
This ensures perfect alignment and prevents layout overcrowding.
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─────────────────────────────────────────────
# MINIMALIST COLOR THEME
# ─────────────────────────────────────────────

COLOR_BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # Pure white
COLOR_TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)      # Deep grey
COLOR_ACCENT = RGBColor(0x00, 0x56, 0xB3)         # Professional blue
COLOR_SUBTEXT = RGBColor(0x7F, 0x8C, 0x8D)        # Soft silver grey


def _set_slide_bg(slide, color: RGBColor):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                  font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────
# MARKDOWN PARSER
# ─────────────────────────────────────────────

def _parse_markdown_slides(structured_text: str, fallback_title: str) -> list[dict]:
    """
    Parse markdown into a list of slide sections.
    """
    slides = []
    current_heading = fallback_title
    current_blocks = []

    lines = structured_text.strip().split("\n")

    for line in lines:
        raw = line.strip()
        if not raw:
            continue

        if raw.startswith("# "):
            if current_blocks or (current_heading and current_heading != fallback_title):
                slides.append({"heading": current_heading, "content_blocks": current_blocks})
            current_heading = raw[2:].strip()
            current_blocks = []

        elif raw.startswith("![IMAGE:"):
            img_path = raw[8:-1]
            if os.path.exists(img_path):
                current_blocks.append({"type": "image", "path": img_path})

        elif raw.startswith("## "):
            current_blocks.append({"type": "bullet", "text": raw[3:].strip()})

        elif raw.startswith("- "):
            clean = re.sub(r"\*\*(.*?)\*\*", r"\1", raw[2:])
            current_blocks.append({"type": "bullet", "text": clean})

        else:
            clean = re.sub(r"\*\*(.*?)\*\*", r"\1", raw)
            if clean:
                current_blocks.append({"type": "bullet", "text": clean})

    if current_blocks or current_heading:
        slides.append({"heading": current_heading, "content_blocks": current_blocks})

    return slides


# ─────────────────────────────────────────────
# SLIDE BUILDERS (SPECIALIZED)
# ─────────────────────────────────────────────

def _build_title_slide(prs, title: str):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, COLOR_BG_WHITE)

    W = prs.slide_width
    H = prs.slide_height

    # Top accent line
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Title
    _add_textbox(
        slide, title,
        Inches(1), Inches(2.3), W - Inches(2), Inches(1.5),
        font_size=42, bold=True, color=COLOR_TEXT_DARK, align=PP_ALIGN.CENTER
    )

    # Subtitle
    _add_textbox(
        slide, "AI-STUDY PRESENTATION",
        Inches(1), Inches(3.8), W - Inches(2), Inches(0.6),
        font_size=18, bold=False, color=COLOR_SUBTEXT, align=PP_ALIGN.CENTER
    )


def _add_visual_slide(prs, heading, image_path):
    """A high-impact slide focusing ONLY on an image."""
    W = prs.slide_width
    H = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_BG_WHITE)

    # Heading
    _add_textbox(
        slide, heading,
        Inches(0.6), Inches(0.3), W - Inches(1.2), Inches(0.6),
        font_size=24, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.LEFT
    )

    # Centered Image
    if os.path.exists(image_path):
        try:
            # We use up to 75% of H, 80% of W
            max_w = W - Inches(2)
            max_h = H - Inches(1.8)
            img_shape = slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=None, height=None)
            
            # Rescale to fit
            factor = min(max_w / img_shape.width, max_h / img_shape.height)
            img_shape.width = int(img_shape.width * factor)
            img_shape.height = int(img_shape.height * factor)
            
            # Center
            img_shape.left = int((W - img_shape.width) / 2)
            img_shape.top = int(Inches(1.2) + (max_h - img_shape.height) / 2)
        except: pass


def _add_content_slide(prs, heading, bullets):
    """A clean detail slide with dynamic font scaling."""
    W = prs.slide_width
    H = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_BG_WHITE)

    # Heading
    _add_textbox(
        slide, heading,
        Inches(0.6), Inches(0.4), W - Inches(1.2), Inches(0.7),
        font_size=28, bold=True, color=COLOR_ACCENT
    )

    bullet_count = len(bullets)
    if not bullet_count: return

    # Content Area
    current_y = Inches(1.3)
    content_h = H - current_y - Inches(0.5)
    txBox = slide.shapes.add_textbox(Inches(0.8), current_y, W - Inches(1.6), content_h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # DYNAMIC SCALING FOR READABILITY (Max 5 bullets)
    if bullet_count > 4: fs, sp = 24, 12
    elif bullet_count > 2: fs, sp = 26, 16
    else: fs, sp = 30, 20

    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.space_before = Pt(sp)
        p.line_spacing = 1.1
        p.level = 0
        p.margin_left = Inches(0.35)
        p.first_line_indent = Inches(-0.25)
        
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(fs)
        run.font.color.rgb = COLOR_TEXT_DARK


def _distribute_content(prs, heading, blocks):
    """Iterate blocks and assign to specialized slides."""
    MAX_BULLETS = 5
    bullets_buffer = []
    slide_index = 0

    for block in blocks:
        if block['type'] == 'image':
            # Flush existing text first
            while bullets_buffer:
                chunk = bullets_buffer[:MAX_BULLETS]
                current_heading = heading if slide_index == 0 else f"{heading} (Cont.)"
                _add_content_slide(prs, current_heading, chunk)
                bullets_buffer = bullets_buffer[MAX_BULLETS:]
                slide_index += 1
            
            # Dedicated Visual Slide
            current_heading = heading if slide_index == 0 else f"{heading} (Cont.)"
            _add_visual_slide(prs, current_heading, block['path'])
            slide_index += 1
        else:
            bullets_buffer.append(block['text'])
            if len(bullets_buffer) >= MAX_BULLETS:
                current_heading = heading if slide_index == 0 else f"{heading} (Cont.)"
                _add_content_slide(prs, current_heading, bullets_buffer)
                bullets_buffer = []
                slide_index += 1

    if bullets_buffer:
        current_heading = heading if slide_index == 0 else f"{heading} (Cont.)"
        _add_content_slide(prs, current_heading, bullets_buffer)


# ─────────────────────────────────────────────
# MAIN GENERATOR
# ─────────────────────────────────────────────

def generate_ppt(
    structured_text: str,
    images: list[dict] = None,
    title: str = "Presentation",
    output_path: str = "data/output/booklet.pptx"
) -> bool:
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        _build_title_slide(prs, title)

        slides_data = _parse_markdown_slides(structured_text, title)
        for s in slides_data:
            _distribute_content(prs, s["heading"], s["content_blocks"])

        prs.save(output_path)
        return True
    except Exception as e:
        print(f"[PPT Generator] Error: {e}")
        return False
